"""
Sync Orchestrator 화면 정의서 PPTX 생성 스크립트
형식: 화면별 1 슬라이드, 헤더 + 설명 + 관련테이블 + 컬럼매핑
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ── 색상 상수 ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
HEADER_BG = RGBColor(0xD6, 0xE4, 0xF0)       # 연한 파란색
SECTION_BG = RGBColor(0x4A, 0x86, 0xC8)       # 진한 파란색
TABLE_HEADER_BG = RGBColor(0xD6, 0xE4, 0xF0)  # 테이블 헤더
TABLE_ALT_BG = RGBColor(0xF2, 0xF2, 0xF2)     # 교대 행
BORDER_COLOR = RGBColor(0x80, 0x80, 0x80)
LIGHT_YELLOW = RGBColor(0xFF, 0xFF, 0xE0)


def set_cell_border(cell, color=BORDER_COLOR, width=Pt(0.5)):
    """셀 테두리 설정"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = etree.SubElement(tcPr, qn(edge))
        ln.set('w', str(int(width)))
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '%02X%02X%02X' % (color[0] if isinstance(color, tuple) else color.red,
                                                color[1] if isinstance(color, tuple) else color.green,
                                                color[2] if isinstance(color, tuple) else color.blue))


def set_cell(cell, text, font_size=8, bold=False, align=PP_ALIGN.CENTER,
             bg_color=None, font_color=BLACK, v_align=MSO_ANCHOR.MIDDLE):
    """셀 텍스트 및 서식 설정"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.font.name = "맑은 고딕"
    p.alignment = align
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = v_align

    # 여백 최소화
    cell.margin_left = Cm(0.1)
    cell.margin_right = Cm(0.1)
    cell.margin_top = Cm(0.05)
    cell.margin_bottom = Cm(0.05)

    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '%02X%02X%02X' % (bg_color[0] if isinstance(bg_color, tuple) else bg_color.red,
                                                bg_color[1] if isinstance(bg_color, tuple) else bg_color.green,
                                                bg_color[2] if isinstance(bg_color, tuple) else bg_color.blue))
    set_cell_border(cell)


def add_header_table(slide, screen_data, prs):
    """상단 헤더 테이블 (프로젝트명, 시스템명, 화면ID 등)"""
    left = Cm(0.5)
    top = Cm(0.3)
    width = prs.slide_width - Cm(1.0)
    height = Cm(2.4)

    # 3행 x 8열 헤더
    table = slide.shapes.add_table(3, 8, left, top, width, height).table

    # 열 너비 설정
    col_widths = [Cm(2.2), Cm(7.5), Cm(1.8), Cm(8.5), Cm(1.8), Cm(3.0), Cm(1.8), Cm(3.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    # 행 높이
    for row in table.rows:
        row.height = Cm(0.8)

    # 1행: 프로젝트명 / 문서명 / 작성자 / 작업일
    set_cell(table.cell(0, 0), "프로젝트명", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(0, 1), "GIMS Sync Orchestrator", 9, False, PP_ALIGN.LEFT)
    set_cell(table.cell(0, 2), "문서명", 9, True, bg_color=HEADER_BG)

    # 문서명 셀 병합 (0,3)
    set_cell(table.cell(0, 3), "화면 정의서", 14, True, PP_ALIGN.CENTER)

    set_cell(table.cell(0, 4), "작성자", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(0, 5), "Claude", 9)
    set_cell(table.cell(0, 6), "작업일", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(0, 7), datetime.date.today().strftime("%Y-%m-%d"), 9)

    # 2행: 시스템명 / (문서명 연속) / 확인자 / 확인일
    set_cell(table.cell(1, 0), "시스템명", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(1, 1), "동기화 관제 시스템", 9, False, PP_ALIGN.LEFT)
    set_cell(table.cell(1, 2), "", 9)
    set_cell(table.cell(1, 3), "", 9)
    set_cell(table.cell(1, 4), "확인자", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(1, 5), "", 9)
    set_cell(table.cell(1, 6), "확인일", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(1, 7), "", 9)

    # 3행: 화면ID / 화면명 / 버전
    set_cell(table.cell(2, 0), "화면ID", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(2, 1), screen_data["screen_id"], 9, False, PP_ALIGN.LEFT)
    set_cell(table.cell(2, 2), "화면명", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(2, 3), screen_data["screen_name"], 9, False, PP_ALIGN.LEFT)

    # 병합: 버전 라벨+값
    set_cell(table.cell(2, 4), "", 9)
    set_cell(table.cell(2, 5), "", 9)
    set_cell(table.cell(2, 6), "버전", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(2, 7), "1.0", 9)

    # 문서명 셀 병합 (row 0-1, col 2-3)
    table.cell(0, 2).merge(table.cell(1, 2))
    table.cell(0, 3).merge(table.cell(1, 3))
    set_cell(table.cell(0, 2), "문서명", 9, True, bg_color=HEADER_BG)
    set_cell(table.cell(0, 3), "화면 정의서", 14, True, PP_ALIGN.CENTER)


def add_section_label(slide, left, top, width, height, text):
    """섹션 라벨 (설 명, 관련테이블 등)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECTION_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Cm(0.2)
    tf.margin_top = Cm(0.1)


def add_text_box(slide, left, top, width, height, text, font_size=9, bold=False):
    """텍스트 박스 추가"""
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.3)
    tf.margin_top = Cm(0.2)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.name = "맑은 고딕"
        p.font.color.rgb = BLACK
        p.space_after = Pt(2)


def add_description_section(slide, screen_data, prs):
    """설명 섹션"""
    right_x = prs.slide_width - Cm(9.5)
    desc_top = Cm(3.0)

    # "설 명" 라벨
    add_section_label(slide, right_x, desc_top, Cm(9.0), Cm(0.7), "설    명")

    # 설명 내용
    desc_text = screen_data.get("description", "")
    add_text_box(slide, right_x, desc_top + Cm(0.8), Cm(9.0), Cm(4.0), desc_text, 9)


def add_related_tables_section(slide, screen_data, prs):
    """관련테이블 섹션"""
    right_x = prs.slide_width - Cm(9.5)
    tbl_top = Cm(7.5)

    # "관련테이블" 라벨
    add_section_label(slide, right_x, tbl_top, Cm(9.0), Cm(0.7), "관련테이블")

    # 테이블 목록
    tables_text = ""
    for i, tbl in enumerate(screen_data.get("related_tables", []), 1):
        tables_text += f"{i}. {tbl['name']}  ({tbl['desc']})\n"

    add_text_box(slide, right_x, tbl_top + Cm(0.8), Cm(9.0), Cm(3.0), tables_text, 9)


def add_screen_description_area(slide, screen_data, prs):
    """화면 설명 영역 (왼쪽 - 화면 구성요소 설명)"""
    left = Cm(0.5)
    top = Cm(3.0)
    width = prs.slide_width - Cm(10.5)
    height = Cm(8.5)

    # 외곽 박스
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(0.5)

    # 화면 경로
    route = screen_data.get("route", "")
    add_text_box(slide, left + Cm(0.3), top + Cm(0.2), width - Cm(0.6), Cm(0.6),
                 f"경로: {route}", 8, True)

    # 화면 구성요소 설명
    elements = screen_data.get("elements", "")
    add_text_box(slide, left + Cm(0.3), top + Cm(1.0), width - Cm(0.6), height - Cm(1.2),
                 elements, 8)


def add_column_mapping_table(slide, screen_data, prs):
    """하단 컬럼 매핑 테이블"""
    tables = screen_data.get("related_tables", [])
    if not tables:
        return

    top = Cm(12.0)
    left = Cm(0.5)
    width = prs.slide_width - Cm(1.0)

    for tbl_data in tables:
        columns = tbl_data.get("columns", [])
        if not columns:
            continue

        num_cols = len(columns)
        if num_cols == 0:
            continue

        # 테이블명 라벨
        add_text_box(slide, left, top, width, Cm(0.5),
                     f"▶ {tbl_data['name']}  ({tbl_data['desc']})", 8, True)
        top += Cm(0.5)

        # 컬럼 테이블: 3행(한글명, 컬럼명, 타입) x N열
        row_count = 3
        tbl = slide.shapes.add_table(row_count, num_cols, left, top, width, Cm(1.8)).table

        col_width = int(width / num_cols)
        for i in range(num_cols):
            tbl.columns[i].width = col_width

        for r in range(row_count):
            tbl.rows[r].height = Cm(0.6)

        for i, col in enumerate(columns):
            set_cell(tbl.cell(0, i), col.get("label", ""), 7, True,
                     bg_color=TABLE_HEADER_BG)
            set_cell(tbl.cell(1, i), col.get("column", ""), 7, False,
                     font_color=RGBColor(0x00, 0x00, 0x80))
            set_cell(tbl.cell(2, i), col.get("type", ""), 7, False,
                     font_color=RGBColor(0x66, 0x66, 0x66))

        top += Cm(2.2)


# ══════════════════════════════════════════════════════
# 화면별 데이터 정의
# ══════════════════════════════════════════════════════

SCREENS = [
    {
        "screen_id": "UI_ORCH_001",
        "screen_name": "대시보드 (Dashboard)",
        "route": "/",
        "description": (
            "1. 실행 통계 카드\n"
            "   - 전체 Agent 수, Online/Offline 수\n"
            "   - 현재 실행 중, 오늘 총 실행, 오늘 실패 건수\n\n"
            "2. Agent 상태 테이블\n"
            "   - Agent명, 코드, Zone, 상태, 마지막 실행 결과\n"
            "   - 카드 클릭 시 필터링 연동\n\n"
            "3. 최근 실행 이력 테이블\n"
            "   - Agent명, 유형, 상태, 읽기/쓰기/스킵 건수\n"
            "   - 소요시간, 트리거 유형, 시작시각\n"
            "   - 10초 주기 자동 갱신"
        ),
        "elements": (
            "[ 통계 카드 영역 ]\n"
            "┌──────┐ ┌──────┐ ┌──────┐ ┌──────┐ ┌──────┐\n"
            "│전체   │ │Online│ │실행중 │ │오늘   │ │오늘   │\n"
            "│Agent │ │Agent │ │      │ │총실행 │ │실패   │\n"
            "└──────┘ └──────┘ └──────┘ └──────┘ └──────┘\n\n"
            "[ Agent 상태 테이블 ]\n"
            "Agent명 | 코드 | Zone | 상태 | 마지막실행\n"
            "──────────────────────────────────────\n\n"
            "[ 최근 실행 이력 테이블 ]\n"
            "Agent | 유형 | 상태 | R/W/S | 소요시간 | 트리거"
        ),
        "related_tables": [
            {
                "name": "agent",
                "desc": "에이전트 등록정보",
                "columns": [
                    {"label": "ID", "column": "id", "type": "BIGINT"},
                    {"label": "에이전트코드", "column": "agent_code", "type": "VARCHAR(50)"},
                    {"label": "에이전트명", "column": "agent_name", "type": "VARCHAR(100)"},
                    {"label": "존", "column": "zone", "type": "VARCHAR(50)"},
                    {"label": "유형", "column": "agent_type", "type": "VARCHAR(20)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "마지막실행", "column": "last_executed_at", "type": "TIMESTAMP"},
                    {"label": "마지막결과", "column": "last_execution_status", "type": "VARCHAR(20)"},
                ]
            },
            {
                "name": "execution_history",
                "desc": "실행 이력",
                "columns": [
                    {"label": "실행ID", "column": "execution_id", "type": "VARCHAR(100)"},
                    {"label": "에이전트코드", "column": "agent_code", "type": "VARCHAR(50)"},
                    {"label": "에이전트명", "column": "agent_name", "type": "VARCHAR(100)"},
                    {"label": "유형", "column": "agent_type", "type": "VARCHAR(30)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "읽기건수", "column": "total_read_count", "type": "BIGINT"},
                    {"label": "쓰기건수", "column": "total_write_count", "type": "BIGINT"},
                    {"label": "스킵건수", "column": "total_skip_count", "type": "BIGINT"},
                    {"label": "소요시간(ms)", "column": "duration_ms", "type": "BIGINT"},
                    {"label": "트리거", "column": "triggered_by", "type": "VARCHAR(20)"},
                    {"label": "시작시각", "column": "started_at", "type": "TIMESTAMP"},
                ]
            },
        ]
    },
    {
        "screen_id": "UI_ORCH_002",
        "screen_name": "Agent 관리 (Agent Management)",
        "route": "/agents",
        "description": (
            "1. Agent 목록 (유형별 그룹)\n"
            "   - RCV(수신), LOADER(적재), SND(송신),\n"
            "     DB_CON_PROXY 유형별 접기/펼치기\n"
            "   - Agent코드, 이름, Zone, 상태, 마지막실행\n\n"
            "2. Agent 등록 (3단계 폼)\n"
            "   Step 1: 엔드포인트 URL 입력 → Agent 탐색\n"
            "   Step 2: Agent 선택, 기본정보 입력\n"
            "   Step 3: 데이터소스/테이블 설정, 스케줄\n\n"
            "3. 액션 버튼\n"
            "   - 헬스체크, 상세보기, 삭제"
        ),
        "elements": (
            "[ Agent 등록 버튼 ]\n\n"
            "▼ RCV (수신) ─────────────────────\n"
            "  코드 | 이름 | Zone | 상태 | 액션\n"
            "  ────────────────────────────────\n\n"
            "▼ LOADER ─────────────────────────\n"
            "  코드 | 이름 | Zone | 상태 | 액션\n\n"
            "▼ SND (송신) ─────────────────────\n"
            "  코드 | 이름 | Zone | 상태 | 액션\n\n"
            "▼ DB_CON_PROXY ───────────────────\n"
            "  코드 | 이름 | Zone | 상태 | 액션"
        ),
        "related_tables": [
            {
                "name": "agent",
                "desc": "에이전트 등록정보",
                "columns": [
                    {"label": "ID", "column": "id", "type": "BIGINT"},
                    {"label": "에이전트코드", "column": "agent_code", "type": "VARCHAR(50)"},
                    {"label": "에이전트명", "column": "agent_name", "type": "VARCHAR(100)"},
                    {"label": "엔드포인트URL", "column": "endpoint_url", "type": "VARCHAR(255)"},
                    {"label": "존", "column": "zone", "type": "VARCHAR(50)"},
                    {"label": "유형", "column": "agent_type", "type": "VARCHAR(20)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "소스DS", "column": "source_datasource_id", "type": "VARCHAR(50)"},
                    {"label": "타겟DS", "column": "target_datasource_id", "type": "VARCHAR(50)"},
                    {"label": "활성화", "column": "is_active", "type": "BOOLEAN"},
                ]
            },
            {
                "name": "agent_table",
                "desc": "에이전트-테이블 매핑",
                "columns": [
                    {"label": "ID", "column": "id", "type": "BIGINT"},
                    {"label": "에이전트ID", "column": "agent_id", "type": "BIGINT (FK)"},
                    {"label": "DS테이블ID", "column": "datasource_table_id", "type": "BIGINT (FK)"},
                    {"label": "테이블유형", "column": "table_type", "type": "VARCHAR(20)"},
                ]
            },
            {
                "name": "schedule",
                "desc": "스케줄 설정 (등록 시 생성)",
                "columns": [
                    {"label": "스케줄ID", "column": "schedule_id", "type": "BIGINT"},
                    {"label": "에이전트ID", "column": "agent_id", "type": "BIGINT (FK)"},
                    {"label": "Cron식", "column": "cron_expression", "type": "VARCHAR(50)"},
                    {"label": "활성화", "column": "is_enabled", "type": "BOOLEAN"},
                ]
            },
        ]
    },
    {
        "screen_id": "UI_ORCH_003",
        "screen_name": "Agent 상세 (Agent Detail)",
        "route": "/agents/[id]",
        "description": (
            "1. 기본정보 탭\n"
            "   - Agent 코드, 이름, 유형, Zone, 상태\n"
            "   - 소스/타겟 데이터소스 연결 정보\n"
            "   - Select Tables (조건실행 대상)\n"
            "   - Retention 설정 (보존일수, 대상 테이블)\n\n"
            "2. 스케줄 관리\n"
            "   - Cron 표현식, 활성/비활성 토글\n"
            "   - 조건실행 옵션 (WHERE 조건 빌더)\n\n"
            "3. 이력 탭\n"
            "   - 해당 Agent 실행 이력 (페이징)\n\n"
            "4. 액션: 헬스체크, 수동실행, 삭제"
        ),
        "elements": (
            "[ Agent명 ]  [상태 뱃지]  [헬스체크] [실행] [삭제]\n"
            "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
            "[ 기본정보 탭 | 이력 탭 ]\n\n"
            "┌ 기본정보 ────────────────────────┐\n"
            "│ 코드: xxx  유형: RCV  Zone: DMZ  │\n"
            "│ 소스DS: xxx  타겟DS: xxx         │\n"
            "│ Retention: 30일                  │\n"
            "└──────────────────────────────────┘\n\n"
            "┌ 조건실행 ────────────────────────┐\n"
            "│ 테이블 | 컬럼 | 연산자 | 값      │\n"
            "└──────────────────────────────────┘\n\n"
            "┌ 스케줄 목록 ────────────────────┐\n"
            "│ Cron | 활성 | 수정 | 삭제        │\n"
            "└──────────────────────────────────┘"
        ),
        "related_tables": [
            {
                "name": "agent",
                "desc": "에이전트 등록정보",
                "columns": [
                    {"label": "에이전트코드", "column": "agent_code", "type": "VARCHAR(50)"},
                    {"label": "에이전트명", "column": "agent_name", "type": "VARCHAR(100)"},
                    {"label": "유형", "column": "agent_type", "type": "VARCHAR(20)"},
                    {"label": "존", "column": "zone", "type": "VARCHAR(50)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "소스DS", "column": "source_datasource_id", "type": "VARCHAR(50)"},
                    {"label": "타겟DS", "column": "target_datasource_id", "type": "VARCHAR(50)"},
                    {"label": "Retention설정", "column": "retention_config", "type": "TEXT (JSON)"},
                ]
            },
            {
                "name": "schedule",
                "desc": "스케줄 설정",
                "columns": [
                    {"label": "스케줄ID", "column": "schedule_id", "type": "BIGINT"},
                    {"label": "에이전트ID", "column": "agent_id", "type": "BIGINT (FK)"},
                    {"label": "Cron식", "column": "cron_expression", "type": "VARCHAR(50)"},
                    {"label": "활성화", "column": "is_enabled", "type": "BOOLEAN"},
                    {"label": "실행옵션", "column": "execution_options", "type": "TEXT (JSON)"},
                    {"label": "생성일", "column": "created_at", "type": "TIMESTAMP"},
                ]
            },
            {
                "name": "execution_history",
                "desc": "실행 이력 (이력 탭)",
                "columns": [
                    {"label": "실행ID", "column": "execution_id", "type": "VARCHAR(100)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "읽기건수", "column": "total_read_count", "type": "BIGINT"},
                    {"label": "쓰기건수", "column": "total_write_count", "type": "BIGINT"},
                    {"label": "소요시간", "column": "duration_ms", "type": "BIGINT"},
                    {"label": "트리거", "column": "triggered_by", "type": "VARCHAR(20)"},
                    {"label": "시작시각", "column": "started_at", "type": "TIMESTAMP"},
                ]
            },
        ]
    },
    {
        "screen_id": "UI_ORCH_004",
        "screen_name": "실행 이력 (Execution History)",
        "route": "/executions",
        "description": (
            "1. 검색 필터 패널\n"
            "   - 상태 (전체/성공/실패/실행중)\n"
            "   - Zone, Agent코드, Agent유형\n"
            "   - 날짜 범위 (시작일 ~ 종료일)\n"
            "   - 텍스트 검색 (Agent명)\n\n"
            "2. 실행 이력 테이블 (페이징)\n"
            "   - Agent명, 코드, 유형, 상태\n"
            "   - 읽기/쓰기/스킵 건수\n"
            "   - 소요시간, 트리거 유형, 시작시각\n"
            "   - 20건 단위 페이징\n"
            "   - 클릭 시 실행 상세로 이동"
        ),
        "elements": (
            "┌ 검색 필터 ──────────────────────┐\n"
            "│ 상태[▼] Zone[▼] Agent코드[▼]    │\n"
            "│ Agent유형[▼]                     │\n"
            "│ 시작일[    ] ~ 종료일[    ]      │\n"
            "│ 검색어[              ] [조회]    │\n"
            "└──────────────────────────────────┘\n\n"
            "총 XX건\n"
            "Agent | 코드 | 유형 | 상태 | R | W | S | 시간 | 트리거\n"
            "──────────────────────────────────────────────────\n"
            "...\n"
            "──────────────────────────────────────────────────\n"
            "         < 1 2 3 4 5 ... >"
        ),
        "related_tables": [
            {
                "name": "execution_history",
                "desc": "실행 이력",
                "columns": [
                    {"label": "실행ID", "column": "execution_id", "type": "VARCHAR(100)"},
                    {"label": "에이전트코드", "column": "agent_code", "type": "VARCHAR(50)"},
                    {"label": "에이전트명", "column": "agent_name", "type": "VARCHAR(100)"},
                    {"label": "유형", "column": "agent_type", "type": "VARCHAR(30)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "읽기건수", "column": "total_read_count", "type": "BIGINT"},
                    {"label": "쓰기건수", "column": "total_write_count", "type": "BIGINT"},
                    {"label": "스킵건수", "column": "total_skip_count", "type": "BIGINT"},
                    {"label": "소요시간(ms)", "column": "duration_ms", "type": "BIGINT"},
                    {"label": "오류메시지", "column": "error_message", "type": "TEXT"},
                    {"label": "시작시각", "column": "started_at", "type": "TIMESTAMP"},
                    {"label": "종료시각", "column": "finished_at", "type": "TIMESTAMP"},
                    {"label": "트리거", "column": "triggered_by", "type": "VARCHAR(20)"},
                ]
            },
            {
                "name": "agent",
                "desc": "에이전트 (필터용)",
                "columns": [
                    {"label": "에이전트코드", "column": "agent_code", "type": "VARCHAR(50)"},
                    {"label": "에이전트명", "column": "agent_name", "type": "VARCHAR(100)"},
                    {"label": "유형", "column": "agent_type", "type": "VARCHAR(20)"},
                    {"label": "존", "column": "zone", "type": "VARCHAR(50)"},
                ]
            },
        ]
    },
    {
        "screen_id": "UI_ORCH_005",
        "screen_name": "실행 상세 (Execution Detail)",
        "route": "/executions/[id]",
        "description": (
            "1. 실행 요약 헤더\n"
            "   - 실행ID, 상태, 시작/종료, 소요시간\n"
            "   - 오류 메시지 (실패 시)\n"
            "   - 총 읽기/쓰기/스킵 건수\n\n"
            "2. 테이블 통계 카드\n"
            "   - 테이블명, 유형(SOURCE/TARGET_IF/TARGET)\n"
            "   - 건수(Total/Success/Failed/Skip)\n\n"
            "3. 데이터 뷰 패널\n"
            "   - 테이블 선택 → 데이터 목록\n"
            "   - 컬럼 정렬, 검색, 상태 필터\n"
            "   - 행 클릭 → Source↔Target 추적\n\n"
            "4. 데이터 추적 (Tracing)\n"
            "   - SOURCE → TARGET 정방향 추적\n"
            "   - TARGET → SOURCE 역방향 추적\n"
            "   - source_refs, PK 기반 매칭"
        ),
        "elements": (
            "┌ 실행 요약 ──────────────────────┐\n"
            "│ ID: xxx-xxx  상태: [SUCCESS]     │\n"
            "│ 시작: 2026-03-20 10:00           │\n"
            "│ 읽기: 150  쓰기: 148  스킵: 2    │\n"
            "└──────────────────────────────────┘\n\n"
            "┌ 테이블 통계 ────────────────────┐\n"
            "│ 테이블명 | 유형 | Total | OK | NG│\n"
            "│ ────────────────────────────────│\n"
            "└──────────────────────────────────┘\n\n"
            "┌ 데이터 뷰 ─────────────────────┐\n"
            "│ [테이블 선택▼] [검색] [필터]     │\n"
            "│ col1 | col2 | col3 | source_refs│\n"
            "│ ── 행 클릭 시 추적 펼침 ──────│\n"
            "└──────────────────────────────────┘"
        ),
        "related_tables": [
            {
                "name": "execution_history",
                "desc": "실행 이력 (요약 정보)",
                "columns": [
                    {"label": "실행ID", "column": "execution_id", "type": "VARCHAR(100)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "읽기건수", "column": "total_read_count", "type": "BIGINT"},
                    {"label": "쓰기건수", "column": "total_write_count", "type": "BIGINT"},
                    {"label": "스킵건수", "column": "total_skip_count", "type": "BIGINT"},
                    {"label": "소요시간", "column": "duration_ms", "type": "BIGINT"},
                    {"label": "오류메시지", "column": "error_message", "type": "TEXT"},
                    {"label": "시작시각", "column": "started_at", "type": "TIMESTAMP"},
                    {"label": "종료시각", "column": "finished_at", "type": "TIMESTAMP"},
                ]
            },
            {
                "name": "execution_step_history",
                "desc": "실행 스텝 이력",
                "columns": [
                    {"label": "ID", "column": "id", "type": "BIGINT"},
                    {"label": "실행ID", "column": "execution_id", "type": "VARCHAR(100)"},
                    {"label": "스텝ID", "column": "step_id", "type": "VARCHAR(100)"},
                    {"label": "상태", "column": "status", "type": "VARCHAR(20)"},
                    {"label": "읽기건수", "column": "read_count", "type": "INTEGER"},
                    {"label": "쓰기건수", "column": "write_count", "type": "INTEGER"},
                    {"label": "스킵건수", "column": "skip_count", "type": "INTEGER"},
                    {"label": "소요시간", "column": "duration_ms", "type": "BIGINT"},
                    {"label": "오류메시지", "column": "error_message", "type": "TEXT"},
                    {"label": "스텝순서", "column": "step_order", "type": "INTEGER"},
                ]
            },
            {
                "name": "datasource / datasource_table",
                "desc": "데이터소스 (테이블명 조회용)",
                "columns": [
                    {"label": "DS ID", "column": "datasource_id", "type": "VARCHAR(50)"},
                    {"label": "DS명", "column": "datasource_name", "type": "VARCHAR(100)"},
                    {"label": "테이블ID", "column": "datasource_table.id", "type": "BIGINT"},
                    {"label": "테이블명", "column": "table_name", "type": "VARCHAR(100)"},
                    {"label": "테이블별칭", "column": "table_alias", "type": "VARCHAR(100)"},
                ]
            },
        ]
    },
    {
        "screen_id": "UI_ORCH_006",
        "screen_name": "DB 관리 (Datasource Management)",
        "route": "/datasources",
        "description": (
            "1. 데이터소스 목록 테이블\n"
            "   - DS ID, 이름, DB유형, Zone\n"
            "   - Host, Port, DB명, 활성상태\n"
            "   - 액션: 테이블관리, 연결테스트, 수정, 삭제\n\n"
            "2. 데이터소스 등록/수정 폼\n"
            "   - 기본정보: ID, 이름, DB유형, 설명\n"
            "   - 연결정보: Host, Port, DB명, 계정, 비밀번호\n"
            "   - Zone 선택, 연결 테스트 버튼\n"
            "   - DB유형 선택 시 기본 포트 자동 입력\n\n"
            "3. 테이블 관리 모달\n"
            "   - 테이블 검색 → 등록 (별칭, 설명)\n"
            "   - 컬럼 메타데이터 자동 수집\n"
            "     (컬럼명, 타입, PK, Nullable, 별칭)\n"
            "   - 등록된 테이블 삭제"
        ),
        "elements": (
            "[ 데이터소스 추가 버튼 ]\n\n"
            "DS ID | 이름 | DB유형 | Zone | Host | Port | DB명 | 상태 | 액션\n"
            "────────────────────────────────────────────────────────────\n"
            "...\n\n"
            "┌ 테이블 관리 모달 ────────────────┐\n"
            "│ [테이블 검색: _______] [검색]     │\n"
            "│                                  │\n"
            "│ ▶ 등록된 테이블 목록              │\n"
            "│ 테이블명 | 별칭 | 설명 | 삭제     │\n"
            "│                                  │\n"
            "│ ▶ 컬럼 상세                      │\n"
            "│ 컬럼명 | 타입 | PK | Null | 별칭  │\n"
            "└──────────────────────────────────┘"
        ),
        "related_tables": [
            {
                "name": "datasource",
                "desc": "데이터소스 연결정보",
                "columns": [
                    {"label": "DS ID", "column": "datasource_id", "type": "VARCHAR(50)"},
                    {"label": "DS명", "column": "datasource_name", "type": "VARCHAR(100)"},
                    {"label": "DB유형", "column": "db_type", "type": "VARCHAR(20)"},
                    {"label": "호스트", "column": "host", "type": "VARCHAR(255)"},
                    {"label": "포트", "column": "port", "type": "INTEGER"},
                    {"label": "DB명", "column": "database_name", "type": "VARCHAR(100)"},
                    {"label": "계정", "column": "username", "type": "VARCHAR(512)"},
                    {"label": "비밀번호", "column": "password", "type": "VARCHAR(1024)"},
                    {"label": "설명", "column": "description", "type": "TEXT"},
                    {"label": "존", "column": "zone", "type": "VARCHAR(50)"},
                    {"label": "활성화", "column": "is_active", "type": "BOOLEAN"},
                ]
            },
            {
                "name": "datasource_table",
                "desc": "등록 테이블 메타데이터",
                "columns": [
                    {"label": "ID", "column": "id", "type": "BIGINT"},
                    {"label": "DS ID", "column": "datasource_id", "type": "VARCHAR(50)"},
                    {"label": "테이블명", "column": "table_name", "type": "VARCHAR(100)"},
                    {"label": "별칭", "column": "table_alias", "type": "VARCHAR(100)"},
                    {"label": "설명", "column": "description", "type": "VARCHAR(500)"},
                    {"label": "생성일", "column": "created_at", "type": "TIMESTAMP"},
                ]
            },
            {
                "name": "datasource_column",
                "desc": "테이블 컬럼 메타데이터",
                "columns": [
                    {"label": "ID", "column": "id", "type": "BIGINT"},
                    {"label": "DS테이블ID", "column": "datasource_table_id", "type": "BIGINT (FK)"},
                    {"label": "컬럼명", "column": "column_name", "type": "VARCHAR(100)"},
                    {"label": "별칭", "column": "column_alias", "type": "VARCHAR(100)"},
                    {"label": "데이터타입", "column": "data_type", "type": "VARCHAR(50)"},
                    {"label": "PK여부", "column": "is_primary_key", "type": "BOOLEAN"},
                    {"label": "Null허용", "column": "is_nullable", "type": "BOOLEAN"},
                    {"label": "설명", "column": "description", "type": "VARCHAR(500)"},
                ]
            },
        ]
    },
]


def main():
    prs = Presentation()
    # 와이드스크린 (16:9)
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    for screen in SCREENS:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

        # 1. 상단 헤더
        add_header_table(slide, screen, prs)

        # 2. 왼쪽: 화면 구성요소 설명
        add_screen_description_area(slide, screen, prs)

        # 3. 오른쪽: 설명 섹션
        add_description_section(slide, screen, prs)

        # 4. 오른쪽: 관련테이블 섹션
        add_related_tables_section(slide, screen, prs)

        # ── 컬럼 매핑은 별도 슬라이드로 ──

    # 컬럼 매핑 슬라이드 (화면별)
    for screen in SCREENS:
        tables = screen.get("related_tables", [])
        if not tables:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 간략 헤더
        add_text_box(slide, Cm(0.5), Cm(0.3), Cm(20), Cm(1.0),
                     f"{screen['screen_id']}  {screen['screen_name']}  —  테이블 컬럼 매핑",
                     12, True)

        top = Cm(1.5)
        left = Cm(0.5)
        width = prs.slide_width - Cm(1.0)

        for tbl_data in tables:
            columns = tbl_data.get("columns", [])
            if not columns:
                continue

            num_cols = len(columns)

            # 테이블명 라벨
            add_text_box(slide, left, top, width, Cm(0.6),
                         f"▶ {tbl_data['name']}  ({tbl_data['desc']})", 9, True)
            top += Cm(0.55)

            # 컬럼 테이블: 3행(한글명, 컬럼명, 타입)
            row_count = 3
            tbl = slide.shapes.add_table(row_count, num_cols, left, top, width, Cm(1.6)).table

            col_width = int(width / num_cols)
            for i in range(num_cols):
                tbl.columns[i].width = col_width

            for r in range(row_count):
                tbl.rows[r].height = Cm(0.5)

            for i, col in enumerate(columns):
                set_cell(tbl.cell(0, i), col.get("label", ""), 7, True,
                         bg_color=TABLE_HEADER_BG)
                set_cell(tbl.cell(1, i), col.get("column", ""), 7, False,
                         font_color=RGBColor(0x00, 0x00, 0x80))
                set_cell(tbl.cell(2, i), col.get("type", ""), 6, False,
                         font_color=RGBColor(0x66, 0x66, 0x66))

            top += Cm(2.0)

            if top > Cm(17):
                # 새 슬라이드 필요
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_text_box(slide, Cm(0.5), Cm(0.3), Cm(20), Cm(1.0),
                             f"{screen['screen_id']}  {screen['screen_name']}  —  테이블 컬럼 매핑 (계속)",
                             12, True)
                top = Cm(1.5)

    output_path = r"D:\dev\claude\GIMS\orchestrator_v2\docs\화면정의서_DB테이블.pptx"
    prs.save(output_path)
    print(f"생성 완료: {output_path}")
    print(f"총 {len(prs.slides)} 슬라이드")


if __name__ == "__main__":
    main()
